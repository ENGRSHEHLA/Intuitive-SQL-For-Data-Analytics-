from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Execution: Turning Plans into Results"
subtitle.text = "Lessons on delegation, accountability, and measurement"

# Helper to add title + bullets slide
def add_bullets(title_text, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1].text_frame
    title.text = title_text
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)

# Slide 2: Main idea
add_bullets("Main Idea", [
    "Managers get results through people",
    "Delegation is the foundation of execution",
])

# Slide 3: Clear Goals & Expectations
add_bullets("Clear Goals & Expectations", [
    "Define what, when, and quality level",
    "Set measurable outcomes and responsibilities",
    "\"Knowing what’s expected\" motivates performance"
])

# Slide 4: Delegation and Execution
add_bullets("Delegation and Execution", [
    "Assign tasks to the right people",
    "Give authority alongside responsibility",
    "Explain tasks, provide resources, allow autonomy",
])

# Slide 5: Right People, Right Roles
add_bullets("Right People in Right Roles", [
    "Match tasks to skills, experience, motivation",
    "Wrong fit → missed deadlines, lower quality",
])

# Slide 6: Measurement, Feedback & Accountability
add_bullets("Measurement, Feedback & Accountability", [
    '\"What gets measured gets done\"',
    "Use metrics, regular feedback, and reviews",
    "Managers must inspect, monitor, and correct",
])

# Slide 7: Develop People & Avoid Reverse Delegation
add_bullets("Develop People & Avoid Reverse Delegation", [
    "Train, coach, and build problem-solvers",
    "Prevent reverse delegation—ask for solutions, not pass tasks back",
    "Good leaders build independent teams that execute",
])

# Save presentation
output_path = "Execution_presentation.pptx"
prs.save(output_path)
print(f"Created {output_path}")
